"""
Generates Vayu_Deck.pptx — import directly into Google Slides.
Run: python generate_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────────────────
BG        = RGBColor(0x06, 0x08, 0x0f)   # deep dark
BLUE      = RGBColor(0x00, 0xc8, 0xff)   # neon blue
PINK      = RGBColor(0xff, 0x3d, 0xa1)   # neon pink
WHITE     = RGBColor(0xff, 0xff, 0xff)
OFFWHITE  = RGBColor(0xe0, 0xe4, 0xf0)
GREY      = RGBColor(0x8a, 0x8f, 0xaa)
DARKCARD  = RGBColor(0x0f, 0x12, 0x21)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def txbox(slide, text, left, top, width, height,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    from pptx.util import Pt, Inches
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def rect(slide, left, top, width, height, fill_color, alpha=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def accent_bar(slide, color=BLUE):
    """Thin gradient accent line at top"""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def slide_number(slide, n):
    txbox(slide, f"{n} / 10", 12.5, 7.1, 0.8, 0.35, size=9, color=GREY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s1 = add_slide()
accent_bar(s1, BLUE)

# Glow card behind title
rect(s1, 3.0, 1.5, 7.33, 4.5, DARKCARD)

txbox(s1, "VAYU", 3.5, 1.8, 6.33, 1.8,
      size=90, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

txbox(s1, "The Autonomous Web Agent", 3.0, 3.5, 7.33, 0.7,
      size=24, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s1, '"The wind that knows everything"', 3.0, 4.2, 7.33, 0.5,
      size=16, italic=True, color=GREY, align=PP_ALIGN.CENTER)

txbox(s1, "Microsoft Build AI Hackathon 2026", 3.0, 5.4, 7.33, 0.4,
      size=13, color=GREY, align=PP_ALIGN.CENTER)

slide_number(s1, 1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
s2 = add_slide()
accent_bar(s2, PINK)

txbox(s2, "THE PROBLEM", 0.5, 0.3, 12.0, 0.5, size=11, color=PINK, bold=True)
txbox(s2, "We spend hours searching.\nAI gives us links. We still do all the work.", 0.5, 0.7, 12.0, 1.2,
      size=30, bold=True, color=WHITE)

pain_points = [
    ("⏱  3–5 hours/week", "lost to web research across dozens of open tabs"),
    ("🔄  Manual cross-referencing", "prices, jobs, courses require comparing multiple sites by hand"),
    ("💬  AI chatbots give summaries", "ChatGPT & Perplexity don't actually browse — they give stale data"),
]

for i, (title, body) in enumerate(pain_points):
    y = 2.2 + i * 1.4
    rect(s2, 0.5, y, 12.33, 1.2, DARKCARD)
    txbox(s2, title, 0.7, y + 0.1, 4.0, 0.5, size=16, bold=True, color=BLUE)
    txbox(s2, body, 0.7, y + 0.55, 11.8, 0.5, size=14, color=OFFWHITE)

txbox(s2, '"The average knowledge worker spends 19% of their week just searching for information"  — McKinsey',
      0.5, 6.9, 12.33, 0.45, size=11, italic=True, color=GREY)

slide_number(s2, 2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════════
s3 = add_slide()
accent_bar(s3, BLUE)

txbox(s3, "THE SOLUTION", 0.5, 0.3, 12.0, 0.5, size=11, color=BLUE, bold=True)
txbox(s3, "Vayu browses, extracts, compares, and\nsynthesises — so you don't have to.",
      0.5, 0.7, 12.0, 1.3, size=30, bold=True, color=WHITE)

txbox(s3, "An autonomous AI agent that navigates real websites, executes multi-step tasks,\nand returns structured answers in plain language.",
      0.5, 2.1, 12.33, 0.8, size=16, color=OFFWHITE)

cols = [
    ("ANY QUERY", BLUE, "Natural language input\nNo special syntax needed\nJust ask like a human"),
    ("ANY WEBSITE", PINK, "Browses real sites live\nFlipkart, Naukri, YouTube\nDevfolio, Codeforces & more"),
    ("STRUCTURED ANSWER", BLUE, "Tables, prices, images\nLinks, recommendations\nClear and actionable"),
]

for i, (title, color, body) in enumerate(cols):
    x = 0.5 + i * 4.28
    rect(s3, x, 3.1, 4.0, 3.5, DARKCARD)
    txbox(s3, title, x + 0.15, 3.3, 3.7, 0.5, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    txbox(s3, body, x + 0.15, 3.9, 3.7, 2.5, size=13, color=OFFWHITE, align=PP_ALIGN.CENTER)

slide_number(s3, 3)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LIVE DEMO RESULTS
# ═══════════════════════════════════════════════════════════════════════════
s4 = add_slide()
accent_bar(s4, PINK)

txbox(s4, "DEMO RESULTS", 0.5, 0.3, 12.0, 0.5, size=11, color=PINK, bold=True)
txbox(s4, "One sentence in. Real results out.", 0.5, 0.7, 12.0, 0.7, size=30, bold=True, color=WHITE)

examples = [
    (BLUE, "CROSS-SITE COMPARISON",
     'User: "Find top 3 budget laptops under ₹40,000 on Flipkart, compare each on Amazon"',
     "→  Comparison table: laptop names, Flipkart price, Amazon price, cheaper site"),
    (PINK, "CONTEXT-AWARE FOLLOW-UP",
     'User: "what are their prices?" (after Python books search)',
     "→  Searches exact book titles from previous message — not random books"),
    (BLUE, "CONDITIONAL LOGIC",
     'User: "If the video has >50k views, find DSA playlist. If not, search GeeksforGeeks."',
     "→  Evaluates condition itself, executes correct branch — no clarification needed"),
]

for i, (color, title, q, r) in enumerate(examples):
    y = 1.65 + i * 1.7
    rect(s4, 0.5, y, 12.33, 1.55, DARKCARD)
    txbox(s4, title, 0.7, y + 0.08, 5.0, 0.4, size=11, bold=True, color=color)
    txbox(s4, q, 0.7, y + 0.45, 11.8, 0.4, size=12, color=GREY, italic=True)
    txbox(s4, r, 0.7, y + 0.9, 11.8, 0.45, size=13, color=OFFWHITE, bold=True)

slide_number(s4, 4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s5 = add_slide()
accent_bar(s5, BLUE)

txbox(s5, "ARCHITECTURE", 0.5, 0.3, 12.0, 0.5, size=11, color=BLUE, bold=True)
txbox(s5, "How Vayu Thinks", 0.5, 0.65, 8.0, 0.7, size=28, bold=True, color=WHITE)

stages = [
    ("USER QUERY", GREY, "Natural language input"),
    ("CONTEXT RESOLVER", BLUE, "Understands follow-ups, names exact items from history"),
    ("GOAL INTERPRETER", PINK, "Detects ambiguity, temporal intent, defines success condition"),
    ("TASK ROUTER", BLUE, "Classifies → Travel | Jobs | Price | Research"),
    ("QUERY ENGINEER", PINK, "Precision Google searches: site:, after:, intitle:, price ranges"),
    ("BROWSER AGENT", BLUE, "Gemini 2.5-flash + browser-use + Playwright on real sites"),
    ("VERIFIER", PINK, "Scores result 0–100, retries with revised plan if quality low"),
    ("SELF-LEARNING MEMORY", BLUE, "Extracts patterns permanently — gets smarter every run"),
]

for i, (name, color, desc) in enumerate(stages):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.5 + row * 1.35
    rect(s5, x, y, 6.1, 1.15, DARKCARD)
    txbox(s5, name, x + 0.15, y + 0.08, 5.8, 0.4, size=12, bold=True, color=color)
    txbox(s5, desc, x + 0.15, y + 0.55, 5.8, 0.45, size=11, color=OFFWHITE)

slide_number(s5, 5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI INTEGRATION
# ═══════════════════════════════════════════════════════════════════════════
s6 = add_slide()
accent_bar(s6, PINK)

txbox(s6, "AI INTEGRATION", 0.5, 0.3, 12.0, 0.5, size=11, color=PINK, bold=True)
txbox(s6, "Built on real AI, not shortcuts.", 0.5, 0.65, 10.0, 0.7, size=28, bold=True, color=WHITE)

ai_blocks = [
    (BLUE, "GEMINI 2.5-FLASH  (Large Model)",
     "Browser agent reasoning, multi-step navigation, conditional logic\nBest-in-class vision + reasoning for web automation"),
    (PINK, "TWO-TIER MODEL ROUTING",
     "Small (gemini-3.1-flash-lite): planning, query engineering, verification\nLarge (gemini-2.5-flash): browser agent — quality matters here"),
    (BLUE, "BROWSER-USE + PLAYWRIGHT",
     "Real Chromium browser — not web scraping\nHandles JS-heavy sites, dropdowns, popups, login walls, lazy loading"),
    (PINK, "SELF-LEARNING MEMORY  (Unique)",
     "Extracts site navigation patterns after every run — stored permanently\nUnlike ChatGPT which forgets everything, Vayu compounds knowledge"),
]

for i, (color, title, body) in enumerate(ai_blocks):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.55 + row * 2.5
    rect(s6, x, y, 6.1, 2.25, DARKCARD)
    txbox(s6, title, x + 0.15, y + 0.12, 5.8, 0.5, size=13, bold=True, color=color)
    txbox(s6, body, x + 0.15, y + 0.65, 5.8, 1.4, size=12, color=OFFWHITE)

slide_number(s6, 6)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DIFFERENTIATION TABLE
# ═══════════════════════════════════════════════════════════════════════════
s7 = add_slide()
accent_bar(s7, BLUE)

txbox(s7, "DIFFERENTIATION", 0.5, 0.3, 12.0, 0.5, size=11, color=BLUE, bold=True)
txbox(s7, "Not a chatbot. Not a search engine.\nSomething entirely new.",
      0.5, 0.65, 10.0, 1.1, size=26, bold=True, color=WHITE)

# Table header
headers = ["Capability", "ChatGPT", "Perplexity", "VAYU"]
header_colors = [WHITE, GREY, GREY, BLUE]
col_widths = [4.5, 2.2, 2.2, 2.5]
col_x = [0.5, 5.1, 7.4, 9.7]

for j, (h, c, w, x) in enumerate(zip(headers, header_colors, col_widths, col_x)):
    align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    rect(s7, x, 1.9, w, 0.4, RGBColor(0x1c, 0x1c, 0x2a))
    txbox(s7, h, x + 0.05, 1.95, w - 0.1, 0.35, size=12, bold=True, color=c, align=align)

rows = [
    ("Browses real websites live",        "❌", "✅", "✅"),
    ("Fills forms, clicks buttons",       "❌", "❌", "✅"),
    ("Learns from experience",            "❌", "❌", "✅"),
    ("Multi-site comparison",             "❌", "❌", "✅"),
    ("Understands follow-up context",     "✅", "❌", "✅"),
    ("Conditional if/else logic",         "❌", "❌", "✅"),
]

for i, row in enumerate(rows):
    y = 2.4 + i * 0.68
    bg = DARKCARD if i % 2 == 0 else RGBColor(0x0a, 0x0d, 0x1a)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_widths)):
        rect(s7, x, y, w, 0.62, bg)
        color = BLUE if j == 3 and cell == "✅" else (WHITE if j == 0 else OFFWHITE)
        bold = j == 3
        align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        txbox(s7, cell, x + 0.05, y + 0.1, w - 0.1, 0.45, size=13, color=color, bold=bold, align=align)

txbox(s7, '"Unlike every other AI tool, Vayu gets measurably smarter with every search it runs."',
      0.5, 6.85, 12.33, 0.45, size=12, italic=True, color=PINK, align=PP_ALIGN.CENTER)

slide_number(s7, 7)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CAPABILITIES
# ═══════════════════════════════════════════════════════════════════════════
s8 = add_slide()
accent_bar(s8, PINK)

txbox(s8, "10 CORE CAPABILITIES", 0.5, 0.3, 12.0, 0.5, size=11, color=PINK, bold=True)
txbox(s8, "Production-grade. All 10. Zero compromises.",
      0.5, 0.65, 10.0, 0.7, size=26, bold=True, color=WHITE)

caps = [
    (BLUE,  "Goal Understanding",        "Ambiguity detection, clarifying questions, success conditions"),
    (PINK,  "Browser Perception",        "DOM + screenshots, structured extraction of prices/tables"),
    (BLUE,  "Smart Navigation",          "Wrong-page detection, redirect handling, 3-step fallback"),
    (PINK,  "Form Interaction",          "Dropdowns, date pickers, modals, popups handled automatically"),
    (BLUE,  "Dynamic Planning",          "Writes plan before acting, replans after 3 stalled steps"),
    (PINK,  "Cross-session Memory",      "Learns site navigation patterns permanently"),
    (BLUE,  "Cross-site Orchestration",  "Chains Flipkart → Amazon → YouTube in one query"),
    (PINK,  "Recovery & Resilience",     "Loop detection, CAPTCHA escape, fallback LLM"),
    (BLUE,  "Self-Verification",         "Scores own result 0–100, retries with revised plan if low"),
    (PINK,  "Conversation Context",      "Resolves follow-ups using exact items from previous messages"),
]

for i, (color, title, desc) in enumerate(caps):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.55 + row * 1.15
    rect(s8, x, y, 6.1, 1.0, DARKCARD)
    txbox(s8, f"✓  {title}", x + 0.15, y + 0.08, 5.8, 0.4, size=13, bold=True, color=color)
    txbox(s8, desc, x + 0.15, y + 0.52, 5.8, 0.38, size=11, color=OFFWHITE)

slide_number(s8, 8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — USE CASES
# ═══════════════════════════════════════════════════════════════════════════
s9 = add_slide()
accent_bar(s9, BLUE)

txbox(s9, "USE CASES", 0.5, 0.3, 12.0, 0.5, size=11, color=BLUE, bold=True)
txbox(s9, "Works for everything. Improves with use.", 0.5, 0.65, 10.0, 0.7, size=26, bold=True, color=WHITE)

use_cases = [
    (BLUE,  "✈️  TRAVEL",       "Find flights Mumbai→Delhi June 15 under ₹8000 with hotel options"),
    (PINK,  "💼  JOBS",          "Python developer jobs in Bangalore salary above 15 LPA on Naukri"),
    (BLUE,  "🛒  PRICE COMPARE", "Compare Samsung S24 price on Flipkart and Amazon India"),
    (PINK,  "🏆  HACKATHONS",    "Top hackathons for BTech CSE engineer skilled in Python and ML"),
    (BLUE,  "📚  LEARNING",      "Best Django course under ₹999 on Udemy + free YouTube playlist"),
    (PINK,  "🔍  RESEARCH",      "Find the latest Hard graph problem on Codeforces and explain it"),
]

for i, (color, title, query) in enumerate(use_cases):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.55 + row * 1.85
    rect(s9, x, y, 6.1, 1.65, DARKCARD)
    txbox(s9, title, x + 0.15, y + 0.12, 5.8, 0.45, size=14, bold=True, color=color)
    txbox(s9, f'"{query}"', x + 0.15, y + 0.65, 5.8, 0.8, size=11, italic=True, color=GREY)

slide_number(s9, 9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — TEAM & LINKS
# ═══════════════════════════════════════════════════════════════════════════
s10 = add_slide()
accent_bar(s10, PINK)

txbox(s10, "TEAM & SUBMISSION", 0.5, 0.3, 12.0, 0.5, size=11, color=PINK, bold=True)
txbox(s10, "Built at Microsoft Build AI Hackathon 2026", 0.5, 0.65, 10.0, 0.7, size=24, bold=True, color=WHITE)

# Team card
rect(s10, 0.5, 1.5, 5.5, 2.2, DARKCARD)
txbox(s10, "Maulik Mahey", 0.7, 1.65, 5.0, 0.55, size=20, bold=True, color=BLUE)
txbox(s10, "Developer & Architect\nmaulik.mahey@indiamart.com", 0.7, 2.2, 5.0, 0.9, size=13, color=OFFWHITE)

# Tech stack
rect(s10, 6.3, 1.5, 6.5, 2.2, DARKCARD)
txbox(s10, "TECH STACK", 6.5, 1.6, 6.0, 0.4, size=11, bold=True, color=PINK)
stack = "Python 3.12  ·  FastAPI  ·  browser-use\nGemini API  ·  Playwright  ·  Railway"
txbox(s10, stack, 6.5, 2.05, 6.0, 1.2, size=13, color=OFFWHITE)

# Links
rect(s10, 0.5, 3.9, 12.33, 2.4, DARKCARD)
txbox(s10, "LINKS", 0.7, 4.0, 11.8, 0.4, size=11, bold=True, color=BLUE)
txbox(s10, "🌐  Live Demo:", 0.7, 4.45, 2.5, 0.4, size=14, color=GREY)
txbox(s10, "your-app.up.railway.app", 3.2, 4.45, 9.0, 0.4, size=14, color=BLUE)
txbox(s10, "💻  GitHub:", 0.7, 4.95, 2.5, 0.4, size=14, color=GREY)
txbox(s10, "github.com/maulik-dot/Microsoft_BuildAI_26", 3.2, 4.95, 9.0, 0.4, size=14, color=BLUE)

txbox(s10, '"The wind moves through all realms unseen — Vayu researches the web the same way."',
      0.5, 6.7, 12.33, 0.55, size=13, italic=True, color=PINK, align=PP_ALIGN.CENTER)

slide_number(s10, 10)

# ═══════════════════════════════════════════════════════════════════════════
prs.save("Vayu_Deck.pptx")
print("✅  Vayu_Deck.pptx created — upload to Google Slides via File → Import slides")
